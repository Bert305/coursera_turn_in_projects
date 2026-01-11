"""
PowerPoint Presentation Generator for SpaceX Coursera Projects
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle=""):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle
    
    return slide

def create_content_slide(prs, title, bullets):
    """Create a slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for bullet in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
    
    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    create_title_slide(prs, 
                      "SpaceX Falcon 9 Landing Prediction",
                      "Data Science Capstone Project\nIBM Data Science Professional Certificate")
    
    # Executive Summary
    create_content_slide(prs, "Executive Summary", [
        "Comprehensive data science project analyzing SpaceX Falcon 9 launches",
        "Objective: Predict first stage landing success to estimate launch costs",
        "Multi-phase approach: Data Collection → Wrangling → Analysis → Visualization → Machine Learning",
        "Key Finding: 83.3% prediction accuracy achieved across multiple ML models",
        "Business Impact: Enables competitive bidding against SpaceX launches",
        "Technologies: Python, SQL, Pandas, Plotly, Dash, Machine Learning (Scikit-learn)"
    ])
    
    # Introduction
    create_content_slide(prs, "Introduction", [
        "SpaceX revolutionized space industry with reusable rocket technology",
        "Falcon 9 launch cost: $62M vs competitors at $165M+",
        "Cost savings primarily from first stage reusability",
        "Project Goal: Build predictive model for first stage landing success",
        "Dataset: Historical SpaceX launch data (2010-present)",
        "7 interconnected analysis modules covering full data science pipeline"
    ])
    
    # Project 1: Data Collection via API
    create_content_slide(prs, "Project 1: Data Collection via SpaceX API", [
        "Utilized SpaceX REST API to collect launch data",
        "Extracted mission details, rocket specifications, payload information",
        "Key data points: Flight number, launch date, booster version, payload mass",
        "Retrieved launch outcomes and landing pad information",
        "Created structured dataset for downstream analysis"
    ])
    
    create_content_slide(prs, "Project 1: Technical Details", [
        "API endpoints: rockets, launches, launchpads, payloads",
        "Data transformation from JSON to Pandas DataFrame",
        "Handled nested JSON structures for complex relationships",
        "Initial dataset: 90+ launches with 17 features",
        "Foundation for all subsequent analysis phases"
    ])
    
    # Project 2: Web Scraping
    create_content_slide(prs, "Project 2: Web Scraping Supplement", [
        "Scraped additional data from Wikipedia and public sources",
        "Complemented API data with historical context",
        "Used BeautifulSoup for HTML parsing",
        "Extracted launch site details and mission descriptions",
        "Enhanced dataset completeness and accuracy"
    ])
    
    create_content_slide(prs, "Project 2: Scraping Methodology", [
        "Targeted Falcon 9 launch records tables",
        "Parsed HTML tables into structured format",
        "Validated data against API sources",
        "Merged with existing dataset for comprehensive view",
        "Ensured data quality through cross-referencing"
    ])
    
    # Project 3: Data Wrangling
    create_content_slide(prs, "Project 3: Data Wrangling & Cleaning", [
        "Identified and handled missing values",
        "LandingPad: 28.57% missing values addressed",
        "PayloadMass: 0% missing - excellent data quality",
        "Created binary classification labels for landing outcome",
        "Standardized categorical variables (Orbit, LaunchSite, etc.)"
    ])
    
    create_content_slide(prs, "Project 3: Feature Engineering", [
        "Converted mission outcomes to binary success/failure",
        "Created one-hot encoded features for categorical variables",
        "Engineered features: Flight history, reuse count, block version",
        "Handled boolean features: GridFins, Reused, Legs",
        "Final dataset: 90 records with 83 engineered features"
    ])
    
    # Project 4: EDA with SQL
    create_content_slide(prs, "Project 4: Exploratory Data Analysis (SQL)", [
        "Utilized SQLite database for structured queries",
        "Analyzed launch patterns across different orbits",
        "Key orbits: GTO (27 launches), ISS (21), VLEO (14)",
        "Investigated success rates by launch site",
        "Examined payload mass relationships with outcomes"
    ])
    
    create_content_slide(prs, "Project 4: SQL Insights", [
        "Most successful launch site: KSC LC-39A",
        "Booster reuse correlation with landing success",
        "Temporal trends: Success rates improved over time",
        "Payload mass range: 500kg to 15,600kg",
        "Mission complexity vs success rate analysis"
    ])
    
    # Project 5: EDA with Visualization
    create_content_slide(prs, "Project 5: Data Visualization", [
        "Created comprehensive visualizations using Matplotlib & Seaborn",
        "Flight number vs payload mass scatter plots",
        "Success rate trends over time",
        "Launch site performance comparisons",
        "Correlation heatmaps for feature relationships"
    ])
    
    create_content_slide(prs, "Project 5: Visual Insights", [
        "Clear improvement in success rates post-2017",
        "CCAFS SLC-40 most frequently used launch site",
        "Payload mass not strongly correlated with landing success",
        "GridFins and Legs usage increased over time",
        "Block 5 boosters showed highest success rates"
    ])
    
    # Project 6: Launch Site Location Analysis
    create_content_slide(prs, "Project 6: Geospatial Analysis", [
        "Mapped launch sites using Folium interactive maps",
        "Launch sites: CCAFS SLC-40, VAFB SLC-4E, KSC LC-39A",
        "Analyzed proximity to coastlines for safety",
        "Visualized success/failure locations on map",
        "Distance calculations between launch sites"
    ])
    
    create_content_slide(prs, "Project 6: Location Insights", [
        "Florida sites (CCAFS, KSC) used for eastward launches",
        "California site (VAFB) for polar orbits",
        "Launch site selection based on orbital requirements",
        "Geographic advantages: Ocean proximity for landing zones",
        "Weather patterns influence success rates"
    ])
    
    # Project 7: Interactive Dashboard
    create_content_slide(prs, "Project 7: Interactive Dash Dashboard", [
        "Built real-time dashboard using Plotly Dash",
        "Interactive dropdowns for launch site filtering",
        "Payload range slider for dynamic filtering",
        "Success rate pie charts by launch site",
        "Scatter plots: Payload mass vs success correlation"
    ])
    
    create_content_slide(prs, "Project 7: Dashboard Features", [
        "Real-time data filtering and visualization",
        "User-friendly interface for stakeholder exploration",
        "Responsive design for multiple screen sizes",
        "Dynamic updates based on user selections",
        "Export capabilities for reports"
    ])
    
    # Project 8: Machine Learning Prediction
    create_content_slide(prs, "Project 8: Machine Learning Models", [
        "Trained 4 classification models: Logistic Regression, SVM, Decision Tree, KNN",
        "Dataset split: 80% training (72 samples), 20% test (18 samples)",
        "Feature standardization using StandardScaler",
        "Hyperparameter tuning with GridSearchCV (10-fold CV)",
        "Model evaluation using accuracy and confusion matrices"
    ])
    
    create_content_slide(prs, "Project 8: Model Results", [
        "Best Model Results (Test Accuracy): 83.3% across all models",
        "Logistic Regression: C=0.01, L2 penalty",
        "SVM: RBF kernel, C=10, gamma=0.001",
        "Decision Tree: Gini criterion, max_depth=12",
        "KNN: n_neighbors=9, algorithm=brute",
        "All models achieved identical test accuracy: 83.33%"
    ])
    
    # Conclusion
    create_content_slide(prs, "Conclusion & Key Findings", [
        "Successfully developed end-to-end data science pipeline",
        "83.3% prediction accuracy demonstrates model reliability",
        "All 4 ML models performed equally well on test data",
        "Data quality and feature engineering were critical success factors",
        "Business value: Competitive intelligence for launch cost estimation"
    ])
    
    create_content_slide(prs, "Recommendations & Future Work", [
        "Collect more data to increase model robustness (currently 90 samples)",
        "Incorporate additional features: Weather, technical specifications",
        "Explore ensemble methods for potentially higher accuracy",
        "Deploy model as API for real-time predictions",
        "Continuous monitoring and model retraining as new launches occur",
        "Expand analysis to include Falcon Heavy and Starship"
    ])
    
    # Thank You Slide
    create_title_slide(prs, 
                      "Thank You",
                      "Questions & Discussion\n\nIBM Data Science Professional Certificate\nCoursera Capstone Project")
    
    # Save presentation
    output_file = "SpaceX_Data_Science_Projects_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created successfully: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
